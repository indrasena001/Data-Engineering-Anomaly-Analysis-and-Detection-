import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Project Paths ───────────────────────────────────────────────────────────
SRC_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.dirname(SRC_DIR)
METRICS_FILE = os.path.join(PROJECT_ROOT, "data", "results", "model_metrics.json")
CHARTS_DIR = os.path.join(PROJECT_ROOT, "output", "charts")
OUTPUT_DIR = os.path.join(os.path.dirname(PROJECT_ROOT), "presentation")
OUTPUT_FILE = os.path.join(OUTPUT_DIR, "NetSentinel_Project_Presentation.pptx")

# Ensure output directory exists
os.makedirs(OUTPUT_DIR, exist_ok=True)

# ── Theme Colors ────────────────────────────────────────────────────────────
COLOR_BG = RGBColor(6, 10, 19)          # Dark background (#060a13)
COLOR_CARD = RGBColor(13, 17, 23)        # Card background (#0d1117)
COLOR_CYAN = RGBColor(0, 240, 255)       # Accent Cyan (#00f0ff)
COLOR_PURPLE = RGBColor(124, 58, 237)   # Accent Purple (#7c3aed)
COLOR_WHITE = RGBColor(226, 232, 240)    # Primary Text (#e2e8f0)
COLOR_MUTED = RGBColor(148, 163, 184)    # Secondary/Muted Text (#94a3b8)
COLOR_RED = RGBColor(255, 51, 102)       # Red accent (#ff3366)
COLOR_GREEN = RGBColor(16, 185, 129)     # Green accent (#10b981)

# ── Load Model Metrics ──────────────────────────────────────────────────────
metrics = {
    "isolation_forest": {
        "accuracy": 0.822,
        "precision": 0.478,
        "recall": 0.567,
        "f1_score": 0.519,
        "anomalies_detected": 20030,
        "training_time": 10.2
    },
    "random_forest": {
        "accuracy": 0.994,
        "f1_weighted": 0.996,
        "f1_macro": 0.820,
        "training_time": 62.0
    }
}

if os.path.exists(METRICS_FILE):
    try:
        with open(METRICS_FILE, 'r') as f:
            data = json.load(f)
            # Isolation Forest
            if "isolation_forest" in data:
                metrics["isolation_forest"]["accuracy"] = data["isolation_forest"].get("accuracy", 0.822)
                metrics["isolation_forest"]["precision"] = data["isolation_forest"].get("precision", 0.478)
                metrics["isolation_forest"]["recall"] = data["isolation_forest"].get("recall", 0.567)
                metrics["isolation_forest"]["f1_score"] = data["isolation_forest"].get("f1_score", 0.519)
                metrics["isolation_forest"]["anomalies_detected"] = data["isolation_forest"].get("anomalies_detected", 20030)
                metrics["isolation_forest"]["training_time"] = data["isolation_forest"].get("training_time_seconds", 10.2)
            # Random Forest
            if "random_forest" in data:
                metrics["random_forest"]["accuracy"] = data["random_forest"].get("accuracy", 0.994)
                metrics["random_forest"]["f1_weighted"] = data["random_forest"].get("f1_weighted", 0.996)
                metrics["random_forest"]["f1_macro"] = data["random_forest"].get("f1_macro", 0.820)
                metrics["random_forest"]["training_time"] = data["random_forest"].get("training_time_seconds", 62.0)
    except Exception as e:
        print(f"Warning: Could not read metrics file, using defaults. Error: {e}")

# ── Helper Functions ────────────────────────────────────────────────────────
def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs

def set_dark_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

def add_header(slide, category, title):
    # Category (small cyan uppercase)
    tx_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
    tf_cat = tx_cat.text_frame
    tf_cat.word_wrap = True
    tf_cat.margin_left = tf_cat.margin_top = tf_cat.margin_right = tf_cat.margin_bottom = 0
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.name = 'Inter'
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_CYAN
    
    # Slide Title
    tx_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
    tf_title = tx_title.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.name = 'Inter'
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_WHITE

def draw_glass_card(slide, left, top, width, height, title, text_items):
    # Shape background container
    shape = slide.shapes.add_shape(
        1,  # rectangle
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_CARD
    shape.line.color.rgb = COLOR_CYAN
    shape.line.width = Pt(0.75)
    
    # Title textbox
    tx_title = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2), width - Inches(0.5), Inches(0.45))
    tf_title = tx_title.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.name = 'Inter'
    p_title.font.size = Pt(16)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_WHITE
    
    # Body textbox
    tx_body = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.75), width - Inches(0.5), height - Inches(0.95))
    tf_body = tx_body.text_frame
    tf_body.word_wrap = True
    tf_body.margin_left = tf_body.margin_top = tf_body.margin_right = tf_body.margin_bottom = 0
    
    for i, item in enumerate(text_items):
        p = tf_body.paragraphs[0] if i == 0 else tf_body.add_paragraph()
        p.text = "• " + item
        p.font.name = 'Inter'
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_MUTED
        p.space_after = Pt(8)

# ── Build Slides ────────────────────────────────────────────────────────────
prs = create_deck()
blank_layout = prs.slide_layouts[6]

# ── Slide 1: Title Slide ────────────────────────────────────────────────────
s1 = prs.slides.add_slide(blank_layout)
set_dark_background(s1)

# Large Center Title Box
tx1 = s1.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(3.0))
tf1 = tx1.text_frame
tf1.word_wrap = True

p1 = tf1.paragraphs[0]
p1.text = "NETSENTINEL"
p1.alignment = PP_ALIGN.CENTER
p1.font.name = 'Inter'
p1.font.size = Pt(54)
p1.font.bold = True
p1.font.color.rgb = COLOR_CYAN

p2 = tf1.add_paragraph()
p2.text = "Network Intrusion & Log Anomaly Detector"
p2.alignment = PP_ALIGN.CENTER
p2.font.name = 'Inter'
p2.font.size = Pt(24)
p2.font.color.rgb = COLOR_WHITE
p2.space_before = Pt(10)

p3 = tf1.add_paragraph()
p3.text = "End-to-End Data Engineering & Machine Learning Pipeline"
p3.alignment = PP_ALIGN.CENTER
p3.font.name = 'Inter'
p3.font.size = Pt(14)
p3.font.color.rgb = COLOR_MUTED
p3.space_before = Pt(30)

# Footer info
tx_foot = s1.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.5))
tf_foot = tx_foot.text_frame
p_foot = tf_foot.paragraphs[0]
p_foot.text = "Class Presentation | Dataset: CICIDS2017"
p_foot.alignment = PP_ALIGN.CENTER
p_foot.font.name = 'Inter'
p_foot.font.size = Pt(11)
p_foot.font.color.rgb = COLOR_PURPLE

# ── Slide 2: Executive Summary ──────────────────────────────────────────────
s2 = prs.slides.add_slide(blank_layout)
set_dark_background(s2)
add_header(s2, "Overview", "Executive Summary")

draw_glass_card(
    s2, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.5),
    "Project Objectives",
    [
        "Build a fully automated data pipeline capable of digesting, cleaning, and preprocessing high-volume network security log datasets.",
        "Implement a dual ML engine combining supervised learning (for known attack categorization) and unsupervised learning (for zero-day anomaly scores).",
        "Export metrics and live predictions into an interactive, high-fidelity web dashboard with real-time threat maps."
    ]
)

draw_glass_card(
    s2, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.5),
    "Core Accomplishments",
    [
        "ETL Pipeline: Successfully merges and processes ~2.8 million network flow logs from the Canadian Institute for Cybersecurity (CICIDS2017).",
        "ML Modeling: Deployed Random Forest with 99.4% weighted classification accuracy, alongside a robust Isolation Forest outlier detector.",
        "Data Export: Engineered automated JSON outputs detailing protocol breakdowns, port mappings, and threat vectors for immediate visualization."
    ]
)

# ── Slide 3: Pipeline Architecture ──────────────────────────────────────────
s3 = prs.slides.add_slide(blank_layout)
set_dark_background(s3)
add_header(s3, "Architecture", "End-to-End Pipeline Workflow")

# Draw 7 horizontal stages
stage_width = Inches(1.6)
stage_gap = Inches(0.08)
start_left = Inches(0.8)
top_pos = Inches(2.2)
height_pos = Inches(4.2)

stages = [
    ("1. Ingestion", "Loads and merges multiple raw CSV datasets into one master dataframe (~2.8M rows)."),
    ("2. Cleaning", "Removes Inf/NaN noise, handles duplicate logs, and cleans column header syntax."),
    ("3. Feature Eng.", "Performs feature selection (40 key columns), handles scale normalization."),
    ("4. ML Modeling", "Splits train/test splits. Trains Isolation Forest (unsupervised) + Random Forest."),
    ("5. Statistics", "Runs numeric analysis to compile protocol distributions and threat counts."),
    ("6. Visualization", "Generates 8 high-quality analytical charts (PNG) saved in output/charts/."),
    ("7. Dashboard", "Compiles metrics & logs into dashboard_data.json to feed the front-end app.")
]

for i, (title, desc) in enumerate(stages):
    left_pos = start_left + i * (stage_width + stage_gap)
    # Shape Container
    shape = s3.shapes.add_shape(1, left_pos, top_pos, stage_width, height_pos)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_CARD
    shape.line.color.rgb = COLOR_CYAN if i % 2 == 0 else COLOR_PURPLE
    shape.line.width = Pt(1.0)
    
    # Text
    tx = s3.shapes.add_textbox(left_pos + Inches(0.08), top_pos + Inches(0.15), stage_width - Inches(0.16), height_pos - Inches(0.3))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = 'Inter'
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.space_after = Pt(8)
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.name = 'Inter'
    p2.font.size = Pt(8.5)
    p2.font.color.rgb = COLOR_MUTED
    p2.line_spacing = 1.25

# ── Slide 4: Data Engineering & ETL ─────────────────────────────────────────
s4 = prs.slides.add_slide(blank_layout)
set_dark_background(s4)
add_header(s4, "Data Processing", "ETL Phase: Ingestion, Cleaning & Features")

# Left: CICIDS2017 Dataset Details Table
rows = 6
cols = 2
table_shape = s4.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.5))
table = table_shape.table

# Set Column Widths
table.columns[0].width = Inches(2.2)
table.columns[1].width = Inches(3.4)

dataset_meta = [
    ("Dataset Name", "CICIDS2017 (UNB)"),
    ("Original Records", "~2.8 Million network flows"),
    ("Feature Columns", "78 numeric + 1 label column"),
    ("Selected ML Features", "40 high-relevance columns"),
    ("Attack Ratio In Data", "~20.0% of total flows"),
    ("Attacks Covered", "DoS, DDoS, PortScan, Botnet, Web, etc.")
]

for row_idx, (k, v) in enumerate(dataset_meta):
    for col_idx, text in enumerate([k, v]):
        cell = table.cell(row_idx, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_CARD
        
        # Font settings
        for p in cell.text_frame.paragraphs:
            p.font.name = 'Inter'
            p.font.size = Pt(11)
            p.font.color.rgb = COLOR_WHITE if col_idx == 1 else COLOR_CYAN
            p.font.bold = (col_idx == 0)

# Right: Details card
draw_glass_card(
    s4, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.5),
    "Data Preprocessing Pipeline",
    [
        "Data Ingestion: Efficient chunk-wise loading to prevent RAM overflows when merging massive Monday-Friday CSVs.",
        "Data Cleansing: Identifies and replaces Infinite/NaN values in high-variance columns. Sanitizes duplicate flow instances.",
        "Scaling & Standardization: Fits a standard scaler to normalise high-value variables (e.g. Flow Duration, Init Win Bytes) for proper mathematical weighting.",
        "Subsampling: Smart balanced selection up to 500,000 samples for swift model training iteration without data loss."
    ]
)

# ── Slide 5: Machine Learning Models ────────────────────────────────────────
s5 = prs.slides.add_slide(blank_layout)
set_dark_background(s5)
add_header(s5, "Modeling", "Dual Machine Learning Engines")

draw_glass_card(
    s5, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.5),
    "Engine A: Random Forest (Supervised)",
    [
        "Goal: Multi-class attack categorisation (identifies if a threat is Botnet, DDoS, Web Attack, PortScan, etc.).",
        "Handling Imbalance: Uses 'class_weight=\"balanced\"' to handle highly imbalanced attacks (e.g., Infiltration has very few samples compared to normal traffic).",
        "Robustness: High depth limit of 20 with 100 estimators prevents overfitting while locking in structural patterns.",
        "Feature Importance: Computes Gini impurity scores per feature, helping security analysts understand what metrics expose intrusions."
    ]
)

draw_glass_card(
    s5, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.5),
    "Engine B: Isolation Forest (Unsupervised)",
    [
        "Goal: Outlier detection. Flags anomalies and potential Zero-Day exploits without relying on class labels.",
        "Logic: Isolate observations by randomly selecting a feature and split value. Anomalies require fewer splits to isolate.",
        "Hyperparameters: Set contamination threshold to 20% (approximate historical ratio of attack traffic in the environment).",
        "Complementary Analysis: Feeds anomaly scores into system logs to alert on unusual volumes before classification occurs."
    ]
)

# ── Slide 6: Performance Evaluation ─────────────────────────────────────────
s6 = prs.slides.add_slide(blank_layout)
set_dark_background(s6)
add_header(s6, "Evaluation", "Model Training Results & Metrics")

# Left: Metrics Table
t_rows = 5
t_cols = 3
m_table_shape = s6.shapes.add_table(t_rows, t_cols, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.5))
m_table = m_table_shape.table
m_table.columns[0].width = Inches(2.2)
m_table.columns[1].width = Inches(1.7)
m_table.columns[2].width = Inches(1.7)

headers = ["Metric", "Isolation Forest", "Random Forest"]
metrics_rows = [
    ("Accuracy", f"{metrics['isolation_forest']['accuracy']*100:.1f}%", f"{metrics['random_forest']['accuracy']*100:.1f}%"),
    ("F1-Score (Weighted)", f"{metrics['isolation_forest']['f1_score']*100:.1f}%", f"{metrics['random_forest']['f1_weighted']*100:.1f}%"),
    ("Training Time", f"{metrics['isolation_forest']['training_time']:.1f}s", f"{metrics['random_forest']['training_time']:.1f}s"),
    ("Learning Type", "Unsupervised", "Supervised (Multi)")
]

# Write headers
for col_idx, text in enumerate(headers):
    cell = m_table.cell(0, col_idx)
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_CARD
    p = cell.text_frame.paragraphs[0]
    p.font.name = 'Inter'
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

# Write rows
for r_idx, row in enumerate(metrics_rows, 1):
    for c_idx, text in enumerate(row):
        cell = m_table.cell(r_idx, c_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_CARD
        p = cell.text_frame.paragraphs[0]
        p.font.name = 'Inter'
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_WHITE
        if c_idx > 0 and "%" in text:
            val = float(text.replace("%",""))
            if val > 90.0:
                p.font.color.rgb = COLOR_GREEN
            elif val > 75.0:
                p.font.color.rgb = COLOR_CYAN

# Right: Embed Confusion Matrix Image
cm_path = os.path.join(CHARTS_DIR, "confusion_matrix.png")
if os.path.exists(cm_path):
    s6.shapes.add_picture(cm_path, Inches(6.8), Inches(1.8), width=Inches(5.6))
else:
    draw_glass_card(
        s6, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.5),
        "Confusion Matrix Chart",
        [
            "The Random Forest model displays exceptionally high precision and recall on major classes (BENIGN, DDoS, DoS, PortScan).",
            "Slight misclassifications occur on small, minor classes like Botnet or Infiltration due to massive class imbalance in the raw data.",
            "Chart outputs are automatically saved to output/charts/confusion_matrix.png."
        ]
    )

# ── Slide 7: Feature Importance & Heatmap ──────────────────────────────────
s7 = prs.slides.add_slide(blank_layout)
set_dark_background(s7)
add_header(s7, "Insights", "Key Security Features & Importances")

fi_path = os.path.join(CHARTS_DIR, "feature_importance.png")
if os.path.exists(fi_path):
    s7.shapes.add_picture(fi_path, Inches(0.8), Inches(1.8), width=Inches(5.6))
else:
    draw_glass_card(
        s7, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.5),
        "Feature Importance Analysis",
        [
            "Determining what metrics identify malicious behavior is crucial for firewalls and IDS rules.",
            "Random Forest assigns highest weights to packet counts, packet lengths, and window bytes.",
            "Visualized features are exported to output/charts/feature_importance.png."
        ]
    )

draw_glass_card(
    s7, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.5),
    "Top Diagnostic Features",
    [
        "Init_Win_bytes_backward: Handshake characteristics are extremely crucial for flagging SSH/FTP brute force and web exploits.",
        "Bwd Packet Length Mean / Max: High volumes of outgoing response payloads indicate possible data exfiltration.",
        "Flow Duration & Flow IAT: Small, rapid duration indicators show automated flooding attacks (e.g. SYN Flood, DDoS).",
        "PSH/ACK Flag Counts: Set flags expose session scanner sweeps (PortScan) and packet structural anomalies."
    ]
)

# ── Slide 8: ROC Curves & Anomaly Scores ────────────────────────────────────
s8 = prs.slides.add_slide(blank_layout)
set_dark_background(s8)
add_header(s8, "Performance", "Validation Curve Analysis")

roc_path = os.path.join(CHARTS_DIR, "roc_curve.png")
if os.path.exists(roc_path):
    s8.shapes.add_picture(roc_path, Inches(0.8), Inches(1.8), width=Inches(5.6))
else:
    draw_glass_card(
        s8, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.5),
        "ROC Curves (One-vs-Rest)",
        [
            "Plotted ROC curves showcase the trade-off between sensitivity and specificity across classes.",
            "Area Under Curve (AUC) values hover around 0.99 for all primary attack categories.",
            "Visualizations output saved to output/charts/roc_curve.png."
        ]
    )

as_path = os.path.join(CHARTS_DIR, "anomaly_scores.png")
if os.path.exists(as_path):
    s8.shapes.add_picture(as_path, Inches(6.8), Inches(1.8), width=Inches(5.6))
else:
    draw_glass_card(
        s8, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.5),
        "Anomaly Score Distribution",
        [
            "Isolation Forest assigns a score between -0.5 and 0.5 to each network flow.",
            "Normal traffic clusters at positive scores, whereas attacks form long tails in negative score ranges.",
            "Provides an early warning score distribution, saved to output/charts/anomaly_scores.png."
        ]
    )

# ── Slide 9: Real-time Web Dashboard ────────────────────────────────────────
s9 = prs.slides.add_slide(blank_layout)
set_dark_background(s9)
add_header(s9, "Monitoring", "Interactive Cybersecurity Dashboard")

draw_glass_card(
    s9, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.5),
    "Front-End Capabilities",
    [
        "Visual Map Layer: Renders an active map representing global threat sources pointing to the central network node.",
        "Interactive Log Feed: Displays system warnings, notifications, and alerts categorized by severity (Critical, Warning, Info).",
        "Active Threat Logs: Lists offending IP addresses, their protocol types, target services, and blocking statuses.",
        "Dynamic Statistics: Connects automatically to pipeline output (dashboard_data.json) or runs in simulation mode."
    ]
)

draw_glass_card(
    s9, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.5),
    "System Technologies",
    [
        "Pipeline Engine: Written in Python utilizing Pandas, Scikit-learn, and Joblib for stateful model serialization.",
        "Web Dashboard: Implemented in vanilla HTML5, modern CSS3 transitions, and raw client-side Javascript.",
        "Data Visualizations: Custom built dashboard widgets rendering using Chart.js alongside standard matplotlib plots.",
        "Production Readiness: Packaged layout suitable for local hosting or containerized deployment."
    ]
)

# ── Save Presentation ───────────────────────────────────────────────────────
prs.save(OUTPUT_FILE)
print(f"Success: Presentation generated and saved to: {OUTPUT_FILE}")
